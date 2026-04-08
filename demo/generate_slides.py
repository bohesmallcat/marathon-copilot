#!/usr/bin/env python3
"""Generate Marathon Copilot demo presentation slides (PPTX)."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Color palette ──
C_BG_DARK   = RGBColor(0x1A, 0x1A, 0x2E)   # deep navy
C_BG_MID    = RGBColor(0x16, 0x21, 0x3E)   # dark blue
C_BG_CARD   = RGBColor(0x1F, 0x2B, 0x4D)   # card bg
C_ACCENT    = RGBColor(0x00, 0xD2, 0xFF)   # cyan accent
C_ACCENT2   = RGBColor(0x00, 0xE6, 0x96)   # green accent
C_ACCENT3   = RGBColor(0xFF, 0x6B, 0x6B)   # red/coral accent
C_ACCENT4   = RGBColor(0xFF, 0xD9, 0x3D)   # yellow accent
C_ACCENT5   = RGBColor(0xC0, 0x84, 0xFC)   # purple accent
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY      = RGBColor(0xB0, 0xB0, 0xC0)
C_LIGHT     = RGBColor(0xE0, 0xE0, 0xF0)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# Use blank layout
blank_layout = prs.slide_layouts[6]

# ── Helper functions ──

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, border_color=None, radius=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, left, top, width, height, text, font_size=18,
             color=C_WHITE, bold=False, align=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_multiline(slide, left, top, width, height, lines, font_size=16,
                  color=C_WHITE, line_space=1.3, font_name="Microsoft YaHei", bold=False):
    """lines: list of (text, color, bold) or str"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if isinstance(line, str):
            txt, clr, bld = line, color, bold
        else:
            txt, clr, bld = line[0], line[1], line[2] if len(line) > 2 else bold
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(font_size)
        p.font.color.rgb = clr
        p.font.bold = bld
        p.font.name = font_name
        p.space_after = Pt(font_size * (line_space - 1) + 2)
    return txBox

def add_icon_text(slide, left, top, icon, text, icon_color=C_ACCENT, text_color=C_WHITE, font_size=18):
    """Add icon (emoji/symbol) + text."""
    add_text(slide, left, top, Inches(0.5), Inches(0.4), icon, font_size + 4, icon_color, font_name="Segoe UI Emoji")
    add_text(slide, left + Inches(0.55), top + Inches(0.02), Inches(5.5), Inches(0.4), text, font_size, text_color)


# ════════════════════════════════════════════
# SLIDE 1: Title
# ════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, C_BG_DARK)

# Decorative accent line
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.2), SLIDE_W, Pt(3))
shape.fill.solid()
shape.fill.fore_color.rgb = C_ACCENT
shape.line.fill.background()

add_text(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1),
         "Marathon Copilot", 54, C_WHITE, True, PP_ALIGN.CENTER)
add_text(slide, Inches(1.5), Inches(2.6), Inches(10), Inches(0.6),
         "你的 AI 马拉松配速教练", 28, C_ACCENT, False, PP_ALIGN.CENTER)

add_text(slide, Inches(1.5), Inches(4.0), Inches(10), Inches(0.8),
         "长期规划 · 科学备赛 · 精准配速 · 深度复盘", 20, C_GRAY, False, PP_ALIGN.CENTER)

add_text(slide, Inches(1.5), Inches(5.8), Inches(10), Inches(0.5),
         "从训练周期化到赛后四维分析，覆盖全马备赛全周期", 16, C_GRAY, False, PP_ALIGN.CENTER)


# ════════════════════════════════════════════
# SLIDE 2: The Problem / Pain Points
# ════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, C_BG_DARK)

add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
         "你有没有遇到过这些问题？", 36, C_WHITE, True)

# Pain point cards
cards = [
    ("比赛配速", "\"我到底该跑多快？\"\n赛前纠结半天，结果一出发就跟着大部队冲，后半程崩得怀疑人生", C_ACCENT3),
    ("天气影响", "\"天气预报说 15°C 多云，PB 天气！\"\n结果到了现场湿度 90%，跑到半程心率就飘了", C_ACCENT4),
    ("训练规划", "\"离比赛还有 4 个月，怎么练？\"\n买了训练计划跟不上，跑量不知道怎么加，怕受伤", C_ACCENT),
    ("赛后复盘", "\"跑完就看个总成绩和平均配速\"\n不知道到底哪里出了问题，下次还是同样的错", C_ACCENT2),
]

for i, (title, desc, accent) in enumerate(cards):
    row = i // 2
    col = i % 2
    x = Inches(0.8) + col * Inches(6.2)
    y = Inches(1.5) + row * Inches(2.8)

    card = add_shape(slide, x, y, Inches(5.8), Inches(2.4), C_BG_CARD, accent)
    add_text(slide, x + Inches(0.4), y + Inches(0.25), Inches(5), Inches(0.5),
             title, 22, accent, True)
    add_multiline(slide, x + Inches(0.4), y + Inches(0.85), Inches(5), Inches(1.4),
                  [desc], 15, C_LIGHT, 1.4)


# ════════════════════════════════════════════
# SLIDE 3: Solution Overview — Full Lifecycle
# ════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, C_BG_DARK)

add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
         "覆盖备赛全周期", 36, C_WHITE, True)
add_text(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.5),
         "从长期训练规划到赛后复盘，每个阶段都有对应的 AI 辅助", 18, C_GRAY)

phases = [
    ("赛前 8-52 周", "长期规划", ["VDOT 递进路径", "分期训练阶段", "跑量周期化递进", "女性周期整合"], C_ACCENT5),
    ("赛前 D-30", "目标设定", ["VDOT 跑力计算", "A/B/C 三档目标", "PB 达成概率", "配速策略选择"], C_ACCENT),
    ("赛前 D-7~D-1", "每日简报", ["天气追踪+环境税", "训练+饮食计划", "恢复建议", "每周天气排课"], C_ACCENT2),
    ("比赛日 D-0", "执行方案", ["逐公里配速表", "心率控制区间", "补给时间表", "熔断机制"], C_ACCENT4),
    ("赛后", "深度复盘", ["四维生物力学", "根因链分析", "偏差分解量化", "能力重估"], C_ACCENT3),
]

for i, (phase, title, items, accent) in enumerate(phases):
    x = Inches(0.3) + i * Inches(2.55)
    y = Inches(1.8)

    # Phase card
    card = add_shape(slide, x, y, Inches(2.35), Inches(5.0), C_BG_CARD, accent)

    # Phase label
    label = add_shape(slide, x + Inches(0.15), y + Inches(0.3), Inches(2.05), Inches(0.45), accent)
    add_text(slide, x + Inches(0.15), y + Inches(0.32), Inches(2.05), Inches(0.45),
             phase, 12, C_BG_DARK, True, PP_ALIGN.CENTER)

    # Title
    add_text(slide, x + Inches(0.15), y + Inches(1.0), Inches(2.05), Inches(0.5),
             title, 22, C_WHITE, True, PP_ALIGN.CENTER)

    # Items
    for j, item in enumerate(items):
        add_text(slide, x + Inches(0.2), y + Inches(1.8) + j * Inches(0.65),
                 Inches(2.0), Inches(0.5),
                 f"  {item}", 14, C_LIGHT)

    # Arrow between cards (except last)
    if i < 4:
        ax = x + Inches(2.35) + Inches(0.02)
        ay = y + Inches(2.2)
        add_text(slide, ax, ay, Inches(0.25), Inches(0.5), "›", 28, C_GRAY, True, PP_ALIGN.CENTER, "Arial")


# ════════════════════════════════════════════
# SLIDE 4: Training Periodization
# ════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, C_BG_DARK)

# Section tag
tag = add_shape(slide, Inches(0.8), Inches(0.35), Inches(1.6), Inches(0.4), C_ACCENT5)
add_text(slide, Inches(0.8), Inches(0.37), Inches(1.6), Inches(0.4),
         "场景一", 14, C_BG_DARK, True, PP_ALIGN.CENTER)

add_text(slide, Inches(2.6), Inches(0.3), Inches(9), Inches(0.6),
         "长期训练规划 —— \"离比赛还有 4 个月，怎么练？\"", 32, C_WHITE, True)

# Left: How it works
add_text(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.5),
         "它怎么帮你", 22, C_ACCENT5, True)

features = [
    ("VDOT 递进路径规划", "根据当前跑力和目标成绩，每 16 周 VDOT 提升 0.5-1.5 点"),
    ("5-6 个分期训练阶段", "恢复期 → 有氧基础 → 有氧发展 → 专项 → 巅峰 → 减量"),
    ("跑量安全递进", "10% 周增量上限 + 每 4 周减载恢复；大体重跑者收紧至 8%"),
    ("里程碑测试验证", "穿插 5km/10km/半马测试赛，验证进步是否达预期"),
    ("女性生理周期整合", "5 相位周期模型，预测比赛日相位，调整训练强度和赛日策略"),
]

for i, (title, desc) in enumerate(features):
    y = Inches(1.9) + i * Inches(1.05)
    add_text(slide, Inches(1.0), y, Inches(5.5), Inches(0.4),
             f"▸ {title}", 17, C_WHITE, True)
    add_text(slide, Inches(1.3), y + Inches(0.38), Inches(5.2), Inches(0.4),
             desc, 14, C_GRAY)

# Right: Example data card
card = add_shape(slide, Inches(7.0), Inches(1.3), Inches(5.5), Inches(4.0), C_BG_CARD, C_ACCENT5)
add_text(slide, Inches(7.3), Inches(1.5), Inches(5), Inches(0.5),
         "训练分期示例", 18, C_ACCENT5, True)

period_data = [
    ("第 1-4 周", "恢复期 · 周跑量 30km · E 配速为主"),
    ("第 5-12 周", "有氧基础期 · 周跑量 40→55km"),
    ("第 13-20 周", "有氧发展期 · 加入 T/I 训练"),
    ("第 21-26 周", "专项期 · MP 配速长距离"),
    ("第 27-30 周", "巅峰期 · 最高质量训练"),
    ("第 31-32 周", "减量期 · 跑量递减至 60%"),
]

for i, (period, desc) in enumerate(period_data):
    y = Inches(2.2) + i * Inches(0.5)
    add_text(slide, Inches(7.5), y, Inches(2.0), Inches(0.4), period, 14, C_ACCENT5, True)
    add_text(slide, Inches(9.5), y, Inches(2.8), Inches(0.4), desc, 14, C_WHITE)

# Probability card
prob_card = add_shape(slide, Inches(7.0), Inches(5.6), Inches(5.5), Inches(1.2), C_BG_CARD, C_ACCENT5)
add_text(slide, Inches(7.3), Inches(5.8), Inches(5), Inches(0.4),
         "多因子达成概率", 15, C_ACCENT5, True)
add_text(slide, Inches(7.3), Inches(6.2), Inches(5), Inches(0.5),
         "VDOT 缺口 × 训练年限 × 跑量 × 伤病 × 耦合比 × 体重 × 年龄", 13, C_LIGHT)


# ════════════════════════════════════════════
# SLIDE 5: Weekly Training Plan
# ════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, C_BG_DARK)

tag = add_shape(slide, Inches(0.8), Inches(0.35), Inches(1.6), Inches(0.4), C_ACCENT5)
add_text(slide, Inches(0.8), Inches(0.37), Inches(1.6), Inches(0.4),
         "场景二", 14, C_BG_DARK, True, PP_ALIGN.CENTER)

add_text(slide, Inches(2.6), Inches(0.3), Inches(9), Inches(0.6),
         "每周训练计划 —— \"这周每天练什么？\"", 32, C_WHITE, True)

# Three columns
weekly_tabs = [
    ("天气自适应排课", C_ACCENT, [
        "周一  轻松跑 8km @ E 配速",
        "  天气：晴 12°C · 适宜度 92%",
        "周二  力量训练（跑步经济性）",
        "  弹性髋 + 单腿深蹲 + 核心",
        "周三  节奏跑 5km @ T 配速",
        "  天气：多云 15°C · 适宜度 85%",
        "周四  休息 / 交叉训练",
        "  天气：暴雨 · 自动安排室内",
    ]),
    ("恢复监控决策树", C_ACCENT2, [
        "每日自检清单",
        "  晨脉：±5 次/分 → 正常",
        "  DOMS：0-3 分制评估",
        "  疼痛量表：0-10 分",
        "  睡眠：≥ 7h → 绿灯",
        "",
        "自动决策",
        "  全绿 → 正常训练",
        "  1 项黄 → 降量 20%",
        "  2+ 项红 → 休息日",
    ]),
    ("饮食 & 力量训练", C_ACCENT3, [
        "训练日饮食",
        "  碳水 6-8g/kg · 蛋白质 1.6g/kg",
        "恢复日饮食",
        "  碳水 4-5g/kg · 蛋白质 1.8g/kg",
        "",
        "力量训练分期",
        "  基础期 → 关节稳定",
        "  发展期 → 跑步经济性",
        "  赛前 → 维持即可",
    ]),
]

for i, (title, accent, items) in enumerate(weekly_tabs):
    x = Inches(0.5) + i * Inches(4.2)
    y = Inches(1.3)

    card = add_shape(slide, x, y, Inches(3.95), Inches(5.7), C_BG_CARD, accent)

    # Tab header
    header = add_shape(slide, x, y, Inches(3.95), Inches(0.6), accent)
    add_text(slide, x, y + Inches(0.05), Inches(3.95), Inches(0.5),
             title, 20, C_BG_DARK, True, PP_ALIGN.CENTER)

    yi = y + Inches(0.85)
    for item in items:
        if item == "":
            yi += Inches(0.15)
            continue
        is_indent = item.startswith("  ")
        txt = item.strip()
        clr = C_GRAY if is_indent else C_WHITE
        fs = 14 if is_indent else 15
        bld = not is_indent
        add_text(slide, x + Inches(0.3 if not is_indent else 0.5),
                 yi, Inches(3.4), Inches(0.5), txt, fs, clr, bld)
        yi += Inches(0.45)

# Bottom note
add_text(slide, Inches(0.8), Inches(7.1), Inches(11), Inches(0.4),
         "所有数值根据跑者体重、VDOT 和当前训练阶段自动计算", 14, C_GRAY, False, PP_ALIGN.CENTER)


# ════════════════════════════════════════════
# SLIDE 6: Race Goal Setting
# ════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, C_BG_DARK)

# Section tag
tag = add_shape(slide, Inches(0.8), Inches(0.35), Inches(1.6), Inches(0.4), C_ACCENT)
add_text(slide, Inches(0.8), Inches(0.37), Inches(1.6), Inches(0.4),
         "场景三", 14, C_BG_DARK, True, PP_ALIGN.CENTER)

add_text(slide, Inches(2.6), Inches(0.3), Inches(8), Inches(0.6),
         "赛前目标设定 —— \"我该跑多快？\"", 32, C_WHITE, True)

# Left: How it works
add_text(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.5),
         "它怎么帮你", 22, C_ACCENT, True)

features = [
    ("根据 PB 反推跑力值（VDOT）", "用你最近的比赛成绩，而非训练数据，算出真实跑力"),
    ("多场比赛数据交叉验证", "对比不同比赛的配速策略、掉速率，找出最优跑法"),
    ("生成 A / B / C 三档目标", "激进冲 PB / 稳健达标 / 保守完赛，不把鸡蛋放一个篮子"),
    ("PB 达成概率评估", "综合跑力、训练量、天气、身体状态，给出胜率百分比"),
    ("GPS 距离修正", "手表显示 42.195km 但你实际跑了 42.6km——修正后配速更准"),
]

for i, (title, desc) in enumerate(features):
    y = Inches(1.9) + i * Inches(1.05)
    add_text(slide, Inches(1.0), y, Inches(5.5), Inches(0.4),
             f"▸ {title}", 17, C_WHITE, True)
    add_text(slide, Inches(1.3), y + Inches(0.38), Inches(5.2), Inches(0.4),
             desc, 14, C_GRAY)

# Right: Example data card
card = add_shape(slide, Inches(7.0), Inches(1.3), Inches(5.5), Inches(5.5), C_BG_CARD, C_ACCENT)
add_text(slide, Inches(7.3), Inches(1.5), Inches(5), Inches(0.5),
         "示例：全马目标设定", 18, C_ACCENT, True)

example_lines = [
    ("跑者 VDOT", "56.0（基于近两场全马 PB）"),
    ("A 目标", "Sub 2:40（配速 3'48\"/km）"),
    ("B 目标", "Sub 2:45（配速 3'55\"/km）"),
    ("C 目标", "Sub 2:50（配速 4'02\"/km）"),
    ("PB 胜率", "65%"),
    ("配速策略", "前慢后快（负分段）"),
    ("距离修正", "42.6km（+0.4km）"),
]

for i, (label, value) in enumerate(example_lines):
    y = Inches(2.2) + i * Inches(0.6)
    add_text(slide, Inches(7.5), y, Inches(2.2), Inches(0.4), label, 15, C_GRAY)
    add_text(slide, Inches(9.5), y, Inches(2.8), Inches(0.4), value, 15, C_WHITE, True)


# ════════════════════════════════════════════
# SLIDE 7: Weather Tracking & Env Tax
# ════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, C_BG_DARK)

tag = add_shape(slide, Inches(0.8), Inches(0.35), Inches(1.6), Inches(0.4), C_ACCENT2)
add_text(slide, Inches(0.8), Inches(0.37), Inches(1.6), Inches(0.4),
         "场景四", 14, C_BG_DARK, True, PP_ALIGN.CENTER)

add_text(slide, Inches(2.6), Inches(0.3), Inches(9), Inches(0.6),
         "天气追踪 & 环境税 —— \"天气对我影响多大？\"", 32, C_WHITE, True)

# Environment Tax explanation
add_text(slide, Inches(0.8), Inches(1.2), Inches(5.5), Inches(0.5),
         "核心概念：环境税", 22, C_ACCENT2, True)
add_multiline(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(1.2),
              ["天气会给你的配速\"加税\"。",
               "同样的心率和体感努力程度，",
               "恶劣天气下你会比理想条件慢 X 秒/公里。",
               "这个 X 就是环境税。"],
              15, C_LIGHT, 1.4)

# Tax breakdown
add_text(slide, Inches(0.8), Inches(3.5), Inches(5.5), Inches(0.5),
         "环境税拆解", 20, C_ACCENT2, True)

tax_items = [
    ("气温", "偏离 10-15°C 最佳区间越远，税越高", "+3\"/km"),
    ("湿度", "> 70% 散热困难，心率漂移", "+5\"/km"),
    ("风速", "逆风 > 15km/h 显著影响", "+2\"/km"),
    ("紫外线", "高 UV 加速脱水和体温上升", "+3\"/km"),
    ("体感温度", "综合以上因素的修正值", "=总税"),
]

for i, (factor, desc, val) in enumerate(tax_items):
    y = Inches(4.1) + i * Inches(0.55)
    add_text(slide, Inches(1.0), y, Inches(1.2), Inches(0.4), factor, 15, C_ACCENT2, True)
    add_text(slide, Inches(2.3), y, Inches(3.0), Inches(0.4), desc, 14, C_GRAY)
    add_text(slide, Inches(5.3), y, Inches(1.0), Inches(0.4), val, 14, C_ACCENT4, True)

# Right: Comparison cards
card = add_shape(slide, Inches(7.0), Inches(1.2), Inches(5.5), Inches(2.5), C_BG_CARD, C_ACCENT2)
add_text(slide, Inches(7.3), Inches(1.4), Inches(5), Inches(0.4),
         "良好天气案例（D-2 预报）", 16, C_ACCENT2, True)
good_items = [
    "气温 10°C → 15°C",
    "湿度 69% → 52%",
    "环境税 ≈ 10\"/km",
    "判定：绿灯，执行原计划",
]
for i, item in enumerate(good_items):
    add_text(slide, Inches(7.5), Inches(1.9) + i * Inches(0.45), Inches(4.5), Inches(0.4),
             item, 15, C_WHITE if i < 3 else C_ACCENT2, i == 3)

# Bad weather card
card2 = add_shape(slide, Inches(7.0), Inches(4.1), Inches(5.5), Inches(2.8), C_BG_CARD, C_ACCENT3)
add_text(slide, Inches(7.3), Inches(4.3), Inches(5), Inches(0.4),
         "恶劣天气案例（赛后复盘数据）", 16, C_ACCENT3, True)
bad_items = [
    "赛道隧道段湿度 > 90%",
    "CO\u2082 浓度升高，体感温度 +3~5°C",
    "环境税 ≈ 92\"/km",
    "结果：超出目标 3 分钟+",
    "判定：红灯，应启动 B 方案降速",
]
for i, item in enumerate(bad_items):
    clr = C_WHITE if i < 3 else (C_ACCENT3 if i >= 3 else C_WHITE)
    add_text(slide, Inches(7.5), Inches(4.8) + i * Inches(0.45), Inches(4.5), Inches(0.4),
             item, 15, clr, i >= 3)


# ════════════════════════════════════════════
# SLIDE 8: Daily Training & Diet
# ════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, C_BG_DARK)

tag = add_shape(slide, Inches(0.8), Inches(0.35), Inches(1.6), Inches(0.4), C_ACCENT4)
add_text(slide, Inches(0.8), Inches(0.37), Inches(1.6), Inches(0.4),
         "场景五", 14, C_BG_DARK, True, PP_ALIGN.CENTER)

add_text(slide, Inches(2.6), Inches(0.3), Inches(9), Inches(0.6),
         "每日训练 & 饮食 —— \"赛前一周怎么练、怎么吃？\"", 32, C_WHITE, True)

# Three tabs
tabs = [
    ("训练计划", C_ACCENT, [
        "D-7  轻松跑 8km @ 5'30\"/km",
        "D-6  节奏跑 5km @ 4'40\"/km",
        "D-5  轻松跑 6km @ 5'30\"/km",
        "D-4  休息日",
        "D-3  轻松跑 4km @ 5'30\"/km",
        "D-2  抖腿跑 3km @ 5'00\"/km",
        "D-1  赛前激活 2km + 3×100m",
    ]),
    ("饮食计划", C_ACCENT2, [
        "D-7 ~ D-4  正常饮食",
        "  碳水 5g/kg · 蛋白质 1.5g/kg",
        "D-3  开始糖原超补",
        "  碳水 8g/kg（60kg = 480g）",
        "D-2  继续超补",
        "  碳水 10g/kg（60kg = 600g）",
        "D-1  高碳水 + 低纤维 + 充分补水",
    ]),
    ("恢复建议", C_ACCENT3, [
        "D-7  泡沫轴全身放松 15min",
        "D-6  拉伸 + 冰浴（如有条件）",
        "D-5  轻度按摩 / 散步",
        "D-4  完全休息，睡足 8h",
        "D-3  泡沫轴下肢 10min",
        "D-2  轻度拉伸 + 早睡",
        "D-1  不做任何高强度拉伸",
    ]),
]

for i, (title, accent, items) in enumerate(tabs):
    x = Inches(0.5) + i * Inches(4.2)
    y = Inches(1.3)

    card = add_shape(slide, x, y, Inches(3.95), Inches(5.7), C_BG_CARD, accent)

    # Tab header
    header = add_shape(slide, x, y, Inches(3.95), Inches(0.6), accent)
    add_text(slide, x, y + Inches(0.05), Inches(3.95), Inches(0.5),
             title, 20, C_BG_DARK, True, PP_ALIGN.CENTER)

    for j, item in enumerate(items):
        is_indent = item.startswith("  ")
        txt = item.strip()
        clr = C_GRAY if is_indent else C_WHITE
        fs = 14 if is_indent else 15
        bld = not is_indent
        add_text(slide, x + Inches(0.3 if not is_indent else 0.5),
                 y + Inches(0.85) + j * Inches(0.63),
                 Inches(3.4), Inches(0.5), txt, fs, clr, bld)

# Bottom note
add_text(slide, Inches(0.8), Inches(7.1), Inches(11), Inches(0.4),
         "所有数值根据跑者体重和 VDOT 自动计算，不同跑者看到的计划不一样", 14, C_GRAY, False, PP_ALIGN.CENTER)


# ════════════════════════════════════════════
# SLIDE 9: Race Day Execution
# ════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, C_BG_DARK)

tag = add_shape(slide, Inches(0.8), Inches(0.35), Inches(1.6), Inches(0.4), C_ACCENT4)
add_text(slide, Inches(0.8), Inches(0.37), Inches(1.6), Inches(0.4),
         "场景六", 14, C_BG_DARK, True, PP_ALIGN.CENTER)

add_text(slide, Inches(2.6), Inches(0.3), Inches(9), Inches(0.6),
         "比赛日执行方案 —— \"比赛当天看什么？\"", 32, C_WHITE, True)

# Status overview
add_text(slide, Inches(0.8), Inches(1.2), Inches(4), Inches(0.5),
         "赛前状态速览", 20, C_ACCENT4, True)

status_card = add_shape(slide, Inches(0.8), Inches(1.8), Inches(5), Inches(1.6), C_BG_CARD, C_ACCENT4)
status_items = [
    ("HRV 评估", "三维分析", "基线偏差+变异系数+赛日回升"),
    ("PB 胜率", "65%", "可以冲但留余量"),
    ("环境税", "10\"/km", "绿灯，执行原计划"),
]
for i, (label, val, note) in enumerate(status_items):
    y = Inches(2.0) + i * Inches(0.45)
    add_text(slide, Inches(1.1), y, Inches(1.2), Inches(0.35), label, 14, C_GRAY)
    add_text(slide, Inches(2.3), y, Inches(1.2), Inches(0.35), val, 15, C_ACCENT4, True)
    add_text(slide, Inches(3.6), y, Inches(2.0), Inches(0.35), note, 13, C_GRAY)

# Pace table
add_text(slide, Inches(0.8), Inches(3.7), Inches(4), Inches(0.5),
         "逐段配速表", 20, C_ACCENT4, True)

pace_card = add_shape(slide, Inches(0.8), Inches(4.3), Inches(5), Inches(3.0), C_BG_CARD)
pace_data = [
    ("0 ~ 5 km", "3'52\"/km", "热身入赛，压住兴奋"),
    ("6 ~ 10 km", "3'50\"/km", "进入巡航节奏"),
    ("11 ~ 30 km", "3'48\"/km", "主力巡航段"),
    ("31 ~ 35 km", "3'50\"/km", "预留撞墙余量"),
    ("36 km →", "3'46\"/km", "有余力则提速冲刺"),
]
for i, (seg, pace, note) in enumerate(pace_data):
    y = Inches(4.5) + i * Inches(0.5)
    add_text(slide, Inches(1.1), y, Inches(1.5), Inches(0.35), seg, 14, C_WHITE)
    add_text(slide, Inches(2.6), y, Inches(1.2), Inches(0.35), pace, 15, C_ACCENT4, True)
    add_text(slide, Inches(3.9), y, Inches(1.8), Inches(0.35), note, 13, C_GRAY)

# Right side: Fuse mechanism
add_text(slide, Inches(6.5), Inches(1.2), Inches(6), Inches(0.5),
         "熔断机制：什么时候该降速？", 20, C_ACCENT3, True)

fuse_card = add_shape(slide, Inches(6.5), Inches(1.8), Inches(6), Inches(2.2), C_BG_CARD, C_ACCENT3)
fuse_rules = [
    ("心率 > 170 bpm 持续 3 分钟", "立即降速至 B 方案配速"),
    ("配速偏差 > 8\"/km 持续 2km", "评估体感，考虑切换 B 方案"),
    ("体感温度 > 28°C", "主动降速 10-15\"/km，增加补水"),
]
for i, (condition, action) in enumerate(fuse_rules):
    y = Inches(2.0) + i * Inches(0.6)
    add_text(slide, Inches(6.8), y, Inches(5.5), Inches(0.35), f"⚠  {condition}", 14, C_ACCENT3, True)
    add_text(slide, Inches(7.2), y + Inches(0.3), Inches(5), Inches(0.35), f"→ {action}", 13, C_LIGHT)

add_multiline(slide, Inches(6.5), Inches(4.2), Inches(6), Inches(0.6),
              [("\"熔断不是认输，是科学止损\"", C_ACCENT4, True)], 18)

# Supplement timeline
add_text(slide, Inches(6.5), Inches(5.0), Inches(6), Inches(0.5),
         "补给时间表", 20, C_ACCENT2, True)

supp_card = add_shape(slide, Inches(6.5), Inches(5.6), Inches(6), Inches(1.7), C_BG_CARD, C_ACCENT2)
supp_items = [
    "赛前 30min   能量胶 1 支 + 200ml 水",
    "每 5km        小口补水 100-150ml",
    "10km / 20km   能量胶各 1 支",
    "30km           最后一支能量胶 + 电解质",
]
for i, item in enumerate(supp_items):
    add_text(slide, Inches(6.8), Inches(5.8) + i * Inches(0.38), Inches(5.5), Inches(0.35),
             item, 14, C_LIGHT)


# ════════════════════════════════════════════
# SLIDE 10: Post-Race Review (Enhanced)
# ════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, C_BG_DARK)

tag = add_shape(slide, Inches(0.8), Inches(0.35), Inches(1.6), Inches(0.4), C_ACCENT3)
add_text(slide, Inches(0.8), Inches(0.37), Inches(1.6), Inches(0.4),
         "场景七", 14, C_BG_DARK, True, PP_ALIGN.CENTER)

add_text(slide, Inches(2.6), Inches(0.3), Inches(9), Inches(0.6),
         "赛后深度复盘 —— \"我为什么崩了？\"", 32, C_WHITE, True)

# Left: What it analyzes
add_text(slide, Inches(0.8), Inches(1.2), Inches(5.5), Inches(0.5),
         "教练级深度分析", 22, C_ACCENT3, True)

analysis_items = [
    ("四维生物力学分析", "步幅 × 步频 × 功率 × 配速，逐 5km 关联分析", C_ACCENT3),
    ("步幅衰减三阶段模型", "稳定期 → 加速衰减 → 崩塌期，量化步幅贡献率（> 70%）", C_ACCENT3),
    ("心率-配速脱耦诊断", "区分心率漂移型（湿度）和糖原耗竭型（配速降 HR 不升）", C_ACCENT3),
    ("偏差分解量化", "总偏差拆解到各因子，秒级精度，总和必须吻合", C_WHITE),
    ("跨赛事对比 & 能力重估", "多场比赛横向对比，去除可消除因素后重估净完赛能力", C_WHITE),
]

for i, (title, desc, clr) in enumerate(analysis_items):
    y = Inches(1.8) + i * Inches(1.0)
    add_text(slide, Inches(1.0), y, Inches(5.5), Inches(0.4),
             f"▸ {title}", 17, clr, True)
    add_text(slide, Inches(1.3), y + Inches(0.38), Inches(5.2), Inches(0.4),
             desc, 14, C_GRAY)

# Right: Real example
card = add_shape(slide, Inches(7.0), Inches(1.2), Inches(5.5), Inches(5.8), C_BG_CARD, C_ACCENT3)
add_text(slide, Inches(7.3), Inches(1.4), Inches(5), Inches(0.4),
         "示例：全马复盘", 18, C_ACCENT3, True)

review_data = [
    ("目标成绩", "2:40:00"),
    ("实际成绩", "2:48:30（+8'30\"）"),
    ("", ""),
    ("分段执行", ""),
    ("0-20km", "完美执行，偏差 < 3\""),
    ("25-30km", "开始掉速，慢 15\"/km"),
    ("35-40km", "严重掉速，慢 40\"/km"),
    ("", ""),
    ("根因链诊断", ""),
    ("结构性根因", "C 级周跑量"),
    ("触发因子", "赛间恢复不充分"),
    ("放大器", "步幅衰减 11%"),
    ("边际因子", "补给时机偏迟"),
    ("", ""),
    ("偏差分解", ""),
    ("跑量不足 +300\"", "恢复不充分 +100\""),
    ("步幅崩塌 +90\"", "其他 +20\" = 8'30\""),
]

yi = Inches(1.9)
for label, value in review_data:
    if label == "" and value == "":
        yi += Inches(0.15)
        continue
    if value == "":
        # Section header
        add_text(slide, Inches(7.5), yi, Inches(4.5), Inches(0.35),
                 f"── {label} ──", 13, C_ACCENT3, True)
        yi += Inches(0.38)
        continue
    add_text(slide, Inches(7.5), yi, Inches(1.8), Inches(0.35), label, 14, C_GRAY)
    is_bad = "+" in value or "衰减" in value or "慢" in value or "耗尽" in value or "不充分" in value or "C 级" in value or "偏迟" in value
    add_text(slide, Inches(9.3), yi, Inches(3.0), Inches(0.35),
             value, 14, C_ACCENT3 if is_bad else C_WHITE, is_bad)
    yi += Inches(0.38)


# ════════════════════════════════════════════
# SLIDE 11: Science behind it (Updated)
# ════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, C_BG_DARK)

add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
         "背后的科学原理", 36, C_WHITE, True)
add_text(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.5),
         "每个数字都有运动生理学依据，不是拍脑袋给建议", 18, C_GRAY)

science = [
    ("VDOT 跑力", "Jack Daniels 公式，全球最主流的跑力评估体系", C_ACCENT),
    ("训练周期化", "Daniels/Pfitzinger/Hansons 方法论，10% 周增量 + 4 周减载 + 分期峰值", C_ACCENT5),
    ("环境修正", "温度/湿度/风速量化模型，身高体重（BSA/体重比）个性化修正", C_ACCENT2),
    ("HRV 评估", "三维度分析（基线偏差 + 变异系数 + 赛日回升）→ 竞技状态", C_ACCENT),
    ("四维生物力学", "步幅-步频-功率-配速关联，步幅衰减三阶段模型，步幅贡献率", C_ACCENT3),
    ("HR-配速脱耦", "前后半程配速/HR 变化，区分心率漂移和糖原耗竭两种崩盘模式", C_ACCENT3),
    ("糖原经济学", "基于体重、配速精算糖原储备 vs 消耗 vs 赛中补充，预测撞墙里程", C_ACCENT3),
    ("营养计算", "基于体重的碳水/蛋白质/脂肪目标，赛前糖原超补递增模型", C_ACCENT2),
    ("女性周期", "5 相位周期模型（Bruinvels 2017 / McNulty 2020），强度+营养+风险调节", C_ACCENT5),
    ("天气排课", "7 天天气预报驱动排课，自动避开恶劣天气日，高温日降速/交叉训练", C_ACCENT2),
]

for i, (title, desc, accent) in enumerate(science):
    y = Inches(1.6) + i * Inches(0.58)
    # Accent dot
    dot = add_shape(slide, Inches(0.9), y + Inches(0.08), Inches(0.15), Inches(0.15), accent)
    add_text(slide, Inches(1.3), y, Inches(2.5), Inches(0.4), title, 15, accent, True)
    add_text(slide, Inches(3.8), y, Inches(8.5), Inches(0.4), desc, 14, C_LIGHT)


# ════════════════════════════════════════════
# SLIDE 12: How to Get Started
# ════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, C_BG_DARK)

add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
         "怎么开始用？", 36, C_WHITE, True)

# Option 1: Mini program
card1 = add_shape(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.0), C_BG_CARD, C_ACCENT)
add_text(slide, Inches(1.1), Inches(1.7), Inches(5), Inches(0.5),
         "方式一：微信小程序", 22, C_ACCENT, True)

mp_steps = [
    "1.  打开小程序，创建跑者档案",
    "     填写身高、体重、PB 成绩",
    "",
    "2.  选择目标比赛",
    "     AI 自动生成配速方案",
    "",
    "3.  每日查看训练 & 饮食计划",
    "     赛前 7 天开始推送",
    "",
    "4.  赛后上传数据",
    "     自动生成复盘报告",
]

yi = Inches(2.4)
for line in mp_steps:
    if line == "":
        yi += Inches(0.1)
        continue
    is_sub = line.startswith("     ")
    add_text(slide, Inches(1.3), yi, Inches(4.8), Inches(0.4),
             line.strip(), 14 if is_sub else 16, C_GRAY if is_sub else C_WHITE, not is_sub)
    yi += Inches(0.38)

# Option 2: Send me your data
card2 = add_shape(slide, Inches(7.0), Inches(1.5), Inches(5.5), Inches(5.0), C_BG_CARD, C_ACCENT2)
add_text(slide, Inches(7.3), Inches(1.7), Inches(5), Inches(0.5),
         "方式二：发给我，帮你生成", 22, C_ACCENT2, True)

add_multiline(slide, Inches(7.3), Inches(2.3), Inches(5), Inches(0.8),
              ["目前小程序还在迭代中，", "你可以把数据发给我，我帮你出方案。"],
              15, C_LIGHT, 1.3)

need_items = [
    ("身高 / 体重", "计算环境税和营养目标"),
    ("历史 PB", "最好有最近 2-3 场比赛数据"),
    ("目标比赛", "名称、日期、距离"),
    ("手表数据", "导出分段配速 + 心率"),
    ("HRV 数据", "可选，有则更准"),
    ("生理周期", "女性跑者可选"),
]

add_text(slide, Inches(7.3), Inches(3.3), Inches(5), Inches(0.4),
         "需要你提供：", 16, C_ACCENT2, True)

for i, (item, note) in enumerate(need_items):
    y = Inches(3.8) + i * Inches(0.48)
    add_text(slide, Inches(7.5), y, Inches(2.0), Inches(0.35), f"▸  {item}", 15, C_WHITE, True)
    add_text(slide, Inches(9.6), y, Inches(2.5), Inches(0.35), note, 13, C_GRAY)

# Bottom
add_text(slide, Inches(0.8), Inches(6.8), Inches(11), Inches(0.4),
         "免费 · 开源 · 你的数据不会上传到任何第三方服务器", 16, C_GRAY, False, PP_ALIGN.CENTER)


# ════════════════════════════════════════════
# SLIDE 13: Thank You
# ════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, C_BG_DARK)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.2), SLIDE_W, Pt(3))
shape.fill.solid()
shape.fill.fore_color.rgb = C_ACCENT
shape.line.fill.background()

add_text(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1),
         "Marathon Copilot", 48, C_WHITE, True, PP_ALIGN.CENTER)

add_text(slide, Inches(1.5), Inches(3.6), Inches(10), Inches(0.6),
         "你的下一场 PB，从科学备赛开始", 24, C_ACCENT, False, PP_ALIGN.CENTER)

add_text(slide, Inches(1.5), Inches(5.0), Inches(10), Inches(0.6),
         "Q & A", 32, C_WHITE, True, PP_ALIGN.CENTER)

add_text(slide, Inches(1.5), Inches(6.2), Inches(10), Inches(0.5),
         "长期规划 · 科学备赛 · 精准配速 · 深度复盘", 16, C_GRAY, False, PP_ALIGN.CENTER)


# ── Save ──
out_dir = os.path.dirname(os.path.abspath(__file__))
os.makedirs(out_dir, exist_ok=True)
out_path = os.path.join(out_dir, "Marathon_Copilot_跑团Demo.pptx")
prs.save(out_path)
print(f"Saved to {out_path}")
print(f"Total slides: {len(prs.slides)}")
